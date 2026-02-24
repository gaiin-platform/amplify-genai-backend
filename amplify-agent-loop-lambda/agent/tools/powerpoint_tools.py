"""
PowerPoint generation tools for the Amplify agent loop.

Uses python-pptx to build presentations programmatically. Presentations are held
in-memory via a module-level dict keyed by UUID. When save_presentation is called,
the file is written to the agent's work_directory — the LambdaFileTracker then
automatically uploads it to S3 and makes it available for download.

Tool flow:
    create_presentation → list_slide_layouts → add_slide → populate_placeholder /
    add_text_box / add_image / add_table / set_speaker_notes → save_presentation
"""

import io
import os
import uuid
import logging
import requests
from typing import List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from agent.components.tool import register_tool
from agent.core import ActionContext

logger = logging.getLogger(__name__)

# Standard slide dimensions (13.33" x 7.5")
SLIDE_WIDTH_IN = 13.33
SLIDE_HEIGHT_IN = 7.5

# In-memory presentation store: { presentation_id: Presentation }
_presentations: dict = {}

# ─────────────────────────────────────────────────────────────────────────────
# Internal helpers
# ─────────────────────────────────────────────────────────────────────────────

def _get_prs(presentation_id: str) -> Presentation:
    """Retrieve a presentation by ID or raise a clear error."""
    prs = _presentations.get(presentation_id)
    if prs is None:
        raise ValueError(
            f"Presentation '{presentation_id}' not found. "
            "Call create_presentation first and use the returned presentation_id."
        )
    return prs


def _parse_color(color_str: Optional[str]) -> Optional[RGBColor]:
    """Parse a hex color string like '#FF0000' or 'FF0000' into RGBColor."""
    if not color_str:
        return None
    color_str = color_str.lstrip("#")
    if len(color_str) != 6:
        return None
    try:
        r = int(color_str[0:2], 16)
        g = int(color_str[2:4], 16)
        b = int(color_str[4:6], 16)
        return RGBColor(r, g, b)
    except ValueError:
        return None


def _alignment_from_str(alignment: str):
    """Convert alignment string to PP_ALIGN value."""
    mapping = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    return mapping.get((alignment or "left").lower(), PP_ALIGN.LEFT)


def _get_s3_templates_bucket() -> str:
    bucket = os.getenv("S3_CONSOLIDATION_BUCKET_NAME")
    if not bucket:
        raise ValueError("S3_CONSOLIDATION_BUCKET_NAME environment variable not set")
    return bucket


# ─────────────────────────────────────────────────────────────────────────────
# Presentation management tools
# ─────────────────────────────────────────────────────────────────────────────

@register_tool(tags=["powerpoint"])
def list_available_templates() -> dict:
    """
    List all PowerPoint templates available to use when creating a presentation.

    Templates are managed by admins and stored in S3. Call this before
    create_presentation to show the user their options. If no templates are
    available, a blank presentation can still be created by omitting
    template_name in create_presentation.

    Returns:
        dict with a 'templates' list of template name strings, and a
        'count' of how many are available.
    """
    import boto3
    from botocore.exceptions import ClientError

    bucket = _get_s3_templates_bucket()
    s3 = boto3.client("s3")

    templates = []
    try:
        paginator = s3.get_paginator("list_objects_v2")
        for page in paginator.paginate(Bucket=bucket, Prefix="powerPointTemplates/"):
            for obj in page.get("Contents", []):
                key = obj["Key"]
                # Strip the prefix and only include .pptx/.potx files
                name = key.replace("powerPointTemplates/", "")
                if name and (name.endswith(".pptx") or name.endswith(".potx")):
                    templates.append(name)
    except ClientError as e:
        return {"error": f"Could not list templates: {str(e)}", "templates": []}

    return {
        "templates": templates,
        "count": len(templates),
        "usage": "Pass a template name to create_presentation(template_name=...) to use it.",
    }

@register_tool(
    tags=["powerpoint"],
    status="Creating presentation...",
    resultStatus="Presentation created",
)
def create_presentation(
    template_name: str = None,
    action_context: ActionContext = None,
) -> dict:
    """
    Create a new PowerPoint presentation, optionally from a brand template stored in S3.

    If template_name is provided (e.g. 'vanderbilt_2023'), the template .pptx/.potx file is
    downloaded from the 'powerPointTemplates/' prefix in the consolidation S3 bucket and used
    as the base. This gives you access to the template's layouts, fonts, and color theme.

    If no template_name is given, a blank presentation with default slide size is created.

    Returns a presentation_id to pass to all other PowerPoint tools.

    Args:
        template_name: Optional name of an S3-hosted template (without extension).
                       Tried as .potx then .pptx.

    Returns:
        dict with keys: presentation_id, slide_count, template_used, layouts
    """
    prs = None

    if template_name:
        import boto3
        from botocore.exceptions import ClientError

        s3 = boto3.client("s3")
        bucket = _get_s3_templates_bucket()

        # Try .potx first, then .pptx
        for ext in (".potx", ".pptx"):
            key = f"powerPointTemplates/{template_name}{ext}"
            try:
                response = s3.get_object(Bucket=bucket, Key=key)
                template_bytes = response["Body"].read()
                prs = Presentation(io.BytesIO(template_bytes))
                logger.info("Loaded template from s3://%s/%s", bucket, key)
                break
            except ClientError as e:
                if e.response["Error"]["Code"] in ("NoSuchKey", "404"):
                    continue
                raise

        if prs is None:
            return {
                "error": f"Template '{template_name}' not found in S3 under powerPointTemplates/",
                "tried_keys": [
                    f"powerPointTemplates/{template_name}.potx",
                    f"powerPointTemplates/{template_name}.pptx",
                ],
            }
    else:
        prs = Presentation()
        # Set standard widescreen dimensions
        prs.slide_width = Inches(SLIDE_WIDTH_IN)
        prs.slide_height = Inches(SLIDE_HEIGHT_IN)

    presentation_id = str(uuid.uuid4())
    _presentations[presentation_id] = prs

    layouts = [
        {"index": i, "name": layout.name}
        for i, layout in enumerate(prs.slide_layouts)
    ]

    return {
        "presentation_id": presentation_id,
        "slide_count": len(prs.slides),
        "template_used": template_name,
        "slide_width_inches": round(prs.slide_width / 914400, 2),
        "slide_height_inches": round(prs.slide_height / 914400, 2),
        "layouts": layouts,
    }


@register_tool(tags=["powerpoint"])
def get_presentation_info(presentation_id: str) -> dict:
    """
    Get summary information about an open presentation.

    Returns slide count, dimensions, and all available slide layouts with their
    placeholder details. Use this to orient yourself before adding content.

    Args:
        presentation_id: ID returned by create_presentation.

    Returns:
        dict with slide_count, dimensions, and layouts list.
    """
    prs = _get_prs(presentation_id)

    layouts = []
    for i, layout in enumerate(prs.slide_layouts):
        placeholders = [
            {
                "idx": ph.placeholder_format.idx,
                "name": ph.name,
                "type": str(ph.placeholder_format.type),
            }
            for ph in layout.placeholders
        ]
        layouts.append({"index": i, "name": layout.name, "placeholders": placeholders})

    return {
        "presentation_id": presentation_id,
        "slide_count": len(prs.slides),
        "slide_width_inches": round(prs.slide_width / 914400, 2),
        "slide_height_inches": round(prs.slide_height / 914400, 2),
        "layouts": layouts,
    }


@register_tool(tags=["powerpoint"])
def get_slide_info(presentation_id: str, slide_index: int) -> dict:
    """
    Get detailed information about a specific slide including all shapes,
    placeholders, and their current content.

    Args:
        presentation_id: ID returned by create_presentation.
        slide_index: 0-based index of the slide.

    Returns:
        dict with layout name, placeholders, and shapes on the slide.
    """
    prs = _get_prs(presentation_id)
    if slide_index < 0 or slide_index >= len(prs.slides):
        return {"error": f"slide_index {slide_index} is out of range (0-{len(prs.slides)-1})"}

    slide = prs.slides[slide_index]

    placeholders = []
    for ph in slide.placeholders:
        text = ""
        try:
            text = ph.text_frame.text if ph.has_text_frame else ""
        except Exception:
            pass
        placeholders.append({
            "idx": ph.placeholder_format.idx,
            "name": ph.name,
            "type": str(ph.placeholder_format.type),
            "current_text": text,
        })

    shapes = []
    for shape in slide.shapes:
        if shape.is_placeholder:
            continue
        shapes.append({
            "name": shape.name,
            "shape_type": str(shape.shape_type),
            "left_inches": round(shape.left / 914400, 2),
            "top_inches": round(shape.top / 914400, 2),
            "width_inches": round(shape.width / 914400, 2),
            "height_inches": round(shape.height / 914400, 2),
        })

    return {
        "slide_index": slide_index,
        "layout_name": slide.slide_layout.name,
        "placeholders": placeholders,
        "shapes": shapes,
    }


# ─────────────────────────────────────────────────────────────────────────────
# Slide content tools
# ─────────────────────────────────────────────────────────────────────────────

@register_tool(
    tags=["powerpoint"],
    status="Adding slide...",
    resultStatus="Added slide at index {result}",
)
def add_slide(
    presentation_id: str,
    layout_index: int = 1,
) -> dict:
    """
    Add a new slide to the presentation using the specified layout.

    Use get_presentation_info first to see available layout indices and names.
    Common layouts: 0=Title Slide, 1=Title and Content, 2=Section Header,
    5=Title Only, 6=Blank.

    Args:
        presentation_id: ID returned by create_presentation.
        layout_index: Index of the slide layout to use (default: 1 = Title and Content).

    Returns:
        dict with slide_index (0-based) and available placeholders.
    """
    prs = _get_prs(presentation_id)

    if layout_index < 0 or layout_index >= len(prs.slide_layouts):
        return {
            "error": f"layout_index {layout_index} out of range (0-{len(prs.slide_layouts)-1})"
        }

    layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(layout)
    slide_index = len(prs.slides) - 1

    placeholders = [
        {
            "idx": ph.placeholder_format.idx,
            "name": ph.name,
            "type": str(ph.placeholder_format.type),
        }
        for ph in slide.placeholders
    ]

    return {
        "slide_index": slide_index,
        "layout_used": layout.name,
        "total_slides": len(prs.slides),
        "placeholders": placeholders,
    }


@register_tool(tags=["powerpoint"])
def populate_placeholder(
    presentation_id: str,
    slide_index: int,
    placeholder_idx: int,
    text: str,
    font_size: int = None,
    bold: bool = None,
    font_color: str = None,
    font_name: str = None,
) -> dict:
    """
    Set the text content of a placeholder on a slide.

    Use get_slide_info or the result of add_slide to find placeholder idx values.
    Placeholder idx 0 is almost always the title; idx 1 is the main content area.

    Supports optional formatting: font size (points), bold, color (hex), and font name.
    Note: This replaces ALL text in the placeholder. For multi-paragraph content,
    use add_bullet_points instead.

    Args:
        presentation_id: ID returned by create_presentation.
        slide_index: 0-based slide index.
        placeholder_idx: The idx value of the placeholder (from add_slide or get_slide_info).
        text: Text content to set.
        font_size: Font size in points (optional).
        bold: Whether text should be bold (optional).
        font_color: Hex color string e.g. '#1F3864' (optional).
        font_name: Font name e.g. 'Calibri' (optional).

    Returns:
        dict confirming the update.
    """
    prs = _get_prs(presentation_id)
    if slide_index < 0 or slide_index >= len(prs.slides):
        return {"error": f"slide_index {slide_index} out of range"}

    slide = prs.slides[slide_index]

    target_ph = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == placeholder_idx:
            target_ph = ph
            break

    if target_ph is None:
        available = [ph.placeholder_format.idx for ph in slide.placeholders]
        return {
            "error": f"Placeholder idx {placeholder_idx} not found on slide {slide_index}",
            "available_idx_values": available,
        }

    if not target_ph.has_text_frame:
        return {"error": f"Placeholder idx {placeholder_idx} does not support text"}

    tf = target_ph.text_frame
    tf.clear()
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = text

    font = run.font
    if font_size is not None:
        font.size = Pt(font_size)
    if bold is not None:
        font.bold = bold
    if font_color:
        color = _parse_color(font_color)
        if color:
            font.color.rgb = color
    if font_name:
        font.name = font_name

    return {
        "status": "ok",
        "slide_index": slide_index,
        "placeholder_idx": placeholder_idx,
        "text_set": text,
    }


@register_tool(tags=["powerpoint"])
def add_bullet_points(
    presentation_id: str,
    slide_index: int,
    placeholder_idx: int,
    bullets: list,
    font_size: int = None,
    font_color: str = None,
    font_name: str = None,
) -> dict:
    """
    Populate a content placeholder with a list of bullet points.

    Each item in bullets can be:
    - A string: "Bullet text"
    - A dict: {"text": "Bullet text", "level": 1, "bold": true, "font_size": 18}
      where level 0 = top-level bullet, 1 = sub-bullet, 2 = sub-sub-bullet.

    This clears any existing content in the placeholder first.

    Args:
        presentation_id: ID returned by create_presentation.
        slide_index: 0-based slide index.
        placeholder_idx: The idx of the content placeholder (usually 1).
        bullets: List of bullet items (strings or dicts as described above).
        font_size: Default font size for all bullets (optional, overridden per-item).
        font_color: Default hex color for all bullets (optional).
        font_name: Default font name for all bullets (optional).

    Returns:
        dict confirming bullets added.
    """
    prs = _get_prs(presentation_id)
    if slide_index < 0 or slide_index >= len(prs.slides):
        return {"error": f"slide_index {slide_index} out of range"}

    slide = prs.slides[slide_index]
    target_ph = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == placeholder_idx:
            target_ph = ph
            break

    if target_ph is None:
        available = [ph.placeholder_format.idx for ph in slide.placeholders]
        return {
            "error": f"Placeholder idx {placeholder_idx} not found",
            "available_idx_values": available,
        }

    if not target_ph.has_text_frame:
        return {"error": f"Placeholder idx {placeholder_idx} does not support text"}

    tf = target_ph.text_frame
    tf.clear()

    default_color = _parse_color(font_color)

    for i, bullet in enumerate(bullets):
        if isinstance(bullet, str):
            item_text = bullet
            item_level = 0
            item_bold = None
            item_size = font_size
            item_color = default_color
            item_font = font_name
        else:
            item_text = bullet.get("text", "")
            item_level = bullet.get("level", 0)
            item_bold = bullet.get("bold", None)
            item_size = bullet.get("font_size", font_size)
            item_color = _parse_color(bullet.get("font_color")) or default_color
            item_font = bullet.get("font_name", font_name)

        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()

        para.level = item_level
        run = para.add_run()
        run.text = item_text

        font = run.font
        if item_size is not None:
            font.size = Pt(item_size)
        if item_bold is not None:
            font.bold = item_bold
        if item_color:
            font.color.rgb = item_color
        if item_font:
            font.name = item_font

    return {
        "status": "ok",
        "slide_index": slide_index,
        "placeholder_idx": placeholder_idx,
        "bullets_added": len(bullets),
    }


@register_tool(tags=["powerpoint"])
def add_text_box(
    presentation_id: str,
    slide_index: int,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 18,
    font_name: str = None,
    bold: bool = False,
    font_color: str = None,
    alignment: str = "left",
) -> dict:
    """
    Add a free-form text box at a specific position on a slide.

    Position and size are in inches from the top-left corner of the slide.
    Standard slide is 13.33" wide × 7.5" tall.

    Common positions:
    - Full width title area: left=0.5, top=0.3, width=12.33, height=1.2
    - Left column: left=0.5, top=1.8, width=5.9, height=5.0
    - Right column: left=6.9, top=1.8, width=5.9, height=5.0
    - Footer: left=0.5, top=6.8, width=12.33, height=0.5

    Args:
        presentation_id: ID returned by create_presentation.
        slide_index: 0-based slide index.
        text: Text content for the text box.
        left: Left edge position in inches.
        top: Top edge position in inches.
        width: Width of the text box in inches.
        height: Height of the text box in inches.
        font_size: Font size in points (default: 18).
        font_name: Font name e.g. 'Calibri' (optional).
        bold: Whether text is bold (default: False).
        font_color: Hex color string e.g. '#FFFFFF' (optional).
        alignment: Text alignment: 'left', 'center', 'right', 'justify' (default: 'left').

    Returns:
        dict confirming the text box was added.
    """
    prs = _get_prs(presentation_id)
    if slide_index < 0 or slide_index >= len(prs.slides):
        return {"error": f"slide_index {slide_index} out of range"}

    slide = prs.slides[slide_index]
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = _alignment_from_str(alignment)
    run = para.add_run()
    run.text = text

    font = run.font
    font.size = Pt(font_size)
    font.bold = bold
    if font_name:
        font.name = font_name
    if font_color:
        color = _parse_color(font_color)
        if color:
            font.color.rgb = color

    return {
        "status": "ok",
        "slide_index": slide_index,
        "left": left,
        "top": top,
        "width": width,
        "height": height,
    }


@register_tool(tags=["powerpoint"])
def add_image(
    presentation_id: str,
    slide_index: int,
    image_source: str,
    left: float,
    top: float,
    width: float = None,
    height: float = None,
) -> dict:
    """
    Add an image to a slide.

    image_source can be:
    - A URL starting with http:// or https:// (image will be downloaded)
    - An absolute local file path (e.g. from the agent's work_directory)

    If only width or only height is provided, the image scales proportionally.
    If neither is provided, the image is inserted at its native size.
    Position is in inches from the top-left corner of the slide.

    Args:
        presentation_id: ID returned by create_presentation.
        slide_index: 0-based slide index.
        image_source: URL or local file path to the image.
        left: Left edge in inches.
        top: Top edge in inches.
        width: Width in inches (optional).
        height: Height in inches (optional).

    Returns:
        dict confirming the image was added.
    """
    prs = _get_prs(presentation_id)
    if slide_index < 0 or slide_index >= len(prs.slides):
        return {"error": f"slide_index {slide_index} out of range"}

    slide = prs.slides[slide_index]

    image_data = None
    if image_source.startswith("http://") or image_source.startswith("https://"):
        try:
            resp = requests.get(image_source, timeout=15)
            resp.raise_for_status()
            image_data = io.BytesIO(resp.content)
        except Exception as e:
            return {"error": f"Failed to download image from URL: {str(e)}"}
    else:
        if not os.path.exists(image_source):
            return {"error": f"Image file not found: {image_source}"}
        image_data = image_source

    pic_args = {
        "left": Inches(left),
        "top": Inches(top),
    }
    if width is not None:
        pic_args["width"] = Inches(width)
    if height is not None:
        pic_args["height"] = Inches(height)

    try:
        slide.shapes.add_picture(image_data, **pic_args)
    except Exception as e:
        return {"error": f"Failed to add image: {str(e)}"}

    return {
        "status": "ok",
        "slide_index": slide_index,
        "left": left,
        "top": top,
        "width": width,
        "height": height,
    }


@register_tool(tags=["powerpoint"])
def add_table(
    presentation_id: str,
    slide_index: int,
    data: list,
    left: float,
    top: float,
    width: float,
    height: float,
    has_header: bool = True,
    font_size: int = 12,
    header_font_color: str = None,
    header_bg_color: str = None,
) -> dict:
    """
    Add a table to a slide.

    data is a list of lists representing rows and columns, e.g.:
        [["Name", "Score", "Grade"], ["Alice", "95", "A"], ["Bob", "82", "B"]]

    If has_header is True, the first row is treated as a header and can be
    styled with header_font_color and header_bg_color.

    Args:
        presentation_id: ID returned by create_presentation.
        slide_index: 0-based slide index.
        data: List of rows, each row is a list of cell strings.
        left: Left edge in inches.
        top: Top edge in inches.
        width: Total table width in inches.
        height: Total table height in inches.
        has_header: Whether to style the first row as a header (default: True).
        font_size: Font size for all cells (default: 12).
        header_font_color: Hex color for header text (optional).
        header_bg_color: Hex color for header background (optional).

    Returns:
        dict confirming the table was added.
    """
    prs = _get_prs(presentation_id)
    if slide_index < 0 or slide_index >= len(prs.slides):
        return {"error": f"slide_index {slide_index} out of range"}

    if not data or not data[0]:
        return {"error": "data must be a non-empty list of rows"}

    slide = prs.slides[slide_index]
    rows = len(data)
    cols = max(len(row) for row in data)

    table = slide.shapes.add_table(
        rows, cols,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    ).table

    h_font_color = _parse_color(header_font_color)
    h_bg_color = _parse_color(header_bg_color)

    for r_idx, row_data in enumerate(data):
        for c_idx in range(cols):
            cell = table.cell(r_idx, c_idx)
            cell_text = row_data[c_idx] if c_idx < len(row_data) else ""

            tf = cell.text_frame
            tf.clear()
            para = tf.paragraphs[0]
            run = para.add_run()
            run.text = str(cell_text)
            run.font.size = Pt(font_size)

            if r_idx == 0 and has_header:
                run.font.bold = True
                if h_font_color:
                    run.font.color.rgb = h_font_color
                if h_bg_color:
                    from pptx.oxml.ns import qn
                    from lxml import etree
                    tc = cell._tc
                    tcPr = tc.get_or_add_tcPr()
                    solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
                    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
                    srgbClr.set("val", str(header_bg_color).lstrip("#"))

    return {
        "status": "ok",
        "slide_index": slide_index,
        "rows": rows,
        "columns": cols,
    }


@register_tool(tags=["powerpoint"])
def set_speaker_notes(
    presentation_id: str,
    slide_index: int,
    notes: str,
) -> dict:
    """
    Set the speaker notes for a slide.

    Speaker notes are shown to the presenter during a presentation but not
    to the audience. Useful for adding talking points or context.

    Args:
        presentation_id: ID returned by create_presentation.
        slide_index: 0-based slide index.
        notes: The notes text to set.

    Returns:
        dict confirming notes were set.
    """
    prs = _get_prs(presentation_id)
    if slide_index < 0 or slide_index >= len(prs.slides):
        return {"error": f"slide_index {slide_index} out of range"}

    slide = prs.slides[slide_index]
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.clear()
    tf.paragraphs[0].add_run().text = notes

    return {"status": "ok", "slide_index": slide_index}


# ─────────────────────────────────────────────────────────────────────────────
# Slide structure tools
# ─────────────────────────────────────────────────────────────────────────────

@register_tool(tags=["powerpoint"])
def duplicate_slide(
    presentation_id: str,
    source_slide_index: int,
    insert_after_index: int = None,
) -> dict:
    """
    Duplicate an existing slide, placing the copy after the specified index.

    Useful when multiple slides share the same layout and you want to clone
    a base slide rather than rebuilding the layout each time.

    Args:
        presentation_id: ID returned by create_presentation.
        source_slide_index: 0-based index of the slide to duplicate.
        insert_after_index: 0-based index to insert after. Defaults to end of deck.

    Returns:
        dict with the new slide's index.
    """
    prs = _get_prs(presentation_id)
    total = len(prs.slides)

    if source_slide_index < 0 or source_slide_index >= total:
        return {"error": f"source_slide_index {source_slide_index} out of range"}

    from pptx.oxml.ns import qn
    import copy

    source_slide = prs.slides[source_slide_index]
    template_xml = copy.deepcopy(source_slide._element)

    # Add a blank slide with the same layout, then replace its XML
    layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)
    new_slide._element.getparent().replace(new_slide._element, template_xml)

    new_index = len(prs.slides) - 1

    # Move to correct position if requested
    if insert_after_index is not None and insert_after_index < new_index:
        xml_slides = prs.slides._sldIdLst
        slides_list = list(xml_slides)
        # Move last element to insert_after_index + 1
        item = slides_list[-1]
        xml_slides.remove(item)
        xml_slides.insert(insert_after_index + 1, item)
        new_index = insert_after_index + 1

    return {
        "status": "ok",
        "new_slide_index": new_index,
        "total_slides": len(prs.slides),
    }


@register_tool(tags=["powerpoint"])
def remove_slide(
    presentation_id: str,
    slide_index: int,
) -> dict:
    """
    Remove a slide from the presentation.

    Args:
        presentation_id: ID returned by create_presentation.
        slide_index: 0-based index of the slide to remove.

    Returns:
        dict with updated slide count.
    """
    prs = _get_prs(presentation_id)
    total = len(prs.slides)

    if slide_index < 0 or slide_index >= total:
        return {"error": f"slide_index {slide_index} out of range (0-{total-1})"}

    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    xml_slides.remove(slides_list[slide_index])

    return {
        "status": "ok",
        "removed_index": slide_index,
        "total_slides": len(prs.slides),
    }


# ─────────────────────────────────────────────────────────────────────────────
# Save tool
# ─────────────────────────────────────────────────────────────────────────────

@register_tool(
    tags=["powerpoint", "file"],
    status="Saving presentation as '{filename}'...",
    resultStatus="Presentation saved: {result}",
)
def save_presentation(
    presentation_id: str,
    filename: str = "presentation.pptx",
    action_context: ActionContext = None,
) -> str:
    """
    Save the presentation to a file and make it available for download.

    The file is written to the agent's work_directory. The LambdaFileTracker
    automatically detects new/changed files at the end of the agent loop,
    uploads them to S3, and returns a presigned download URL to the user.

    IMPORTANT: Always call this as the final step after all slides are built.
    The filename should end in .pptx.

    Args:
        presentation_id: ID returned by create_presentation.
        filename: Output filename (default: 'presentation.pptx').
                  Use a descriptive name like 'Q4_report.pptx'.
        action_context: Injected automatically by the agent loop.

    Returns:
        The full file path where the presentation was saved.
    """
    prs = _get_prs(presentation_id)

    work_dir = None
    if action_context:
        work_dir = action_context.get("work_directory")
    if not work_dir:
        work_dir = "/tmp"

    if not filename.lower().endswith(".pptx"):
        filename = filename + ".pptx"

    output_path = os.path.join(work_dir, filename)

    try:
        prs.save(output_path)
        logger.info("Presentation saved to %s", output_path)
    except Exception as e:
        return f"Error saving presentation: {str(e)}"

    # Clean up from memory after saving
    del _presentations[presentation_id]

    return output_path
