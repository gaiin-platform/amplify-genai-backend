# Copyright (c) 2024 Vanderbilt University
# Authors: Jules White, Allen Karns, Karely Rodriguez, Max Moundas

from pptx import Presentation
import io

from rag.handlers.text import TextExtractionHandler


class PPTXHandler(TextExtractionHandler):
    def extract_text(self, file_content, encoding):
        with io.BytesIO(file_content) as f:
            prs = Presentation(f)

            chunks = []
            for slide_number, slide in enumerate(prs.slides, start=1):
                # Extract text from this slide's shapes
                slide_text_parts = [
                    shape.text
                    for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text
                ]
                slide_text = " ".join(slide_text_parts).strip()
                if slide_text:
                    # Create a single chunk per slide with the slide text
                    chunks.append(
                        {
                            "content": slide_text,
                            "location": {"slide_number": slide_number},
                            "canSplit": True,
                        }
                    )

            return chunks

    def handle(self, file_content, key, split_params):
        # Load the presentation
        with io.BytesIO(file_content) as f:
            prs = Presentation(f)

            chunks = []
            for slide_number, slide in enumerate(prs.slides, start=1):
                # Extract text from this slide's shapes
                slide_text_parts = [
                    shape.text
                    for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text
                ]
                slide_text = " ".join(slide_text_parts).strip()
                if slide_text:
                    # Create a single chunk per slide with the slide text
                    chunks.append(
                        {
                            "content": slide_text,
                            "location": {"slide_number": slide_number},
                        }
                    )

            return chunks
