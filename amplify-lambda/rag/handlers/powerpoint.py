# Copyright (c) 2024 Vanderbilt University
# Authors: Jules White, Allen Karns, Karely Rodriguez, Max Moundas

from pptx import Presentation
import io
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw, ImageFont
from rag.handlers.text import TextExtractionHandler
from enum import Enum
from rag.handlers.shared_functions import (
    hash_visual_data,
    ensure_supported_format,
    format_visual_chunk_data,
)
from pptx.util import Pt
from pptx.dml.color import RGBColor

PNG = "image/png"


class VisualType(Enum):
    """Enum for visual content types we extract from PowerPoint"""

    IMAGE = "Image"  # Photos, screenshots, imported images
    DIAGRAM = "Diagram"  # Complex grouped shapes and diagrams
    FREEFORM = "Freeform"  # Custom drawn shapes, annotations, arrows


class PPTXHandler(TextExtractionHandler):
    def extract_text(self, file_content, visual_map={}):
        """
        Legacy method for direct text extraction (used as fallback).
        The main flow now uses preprocess_pptx_visuals + MarkItDown.
        """
        with io.BytesIO(file_content) as f:
            prs = Presentation(f)

            # Preprocess visual_map to group visuals by slide number for efficient lookup
            visuals_by_slide = {}
            for visual_marker, visual_data in visual_map.items():
                slide_num = visual_data.get("location", {}).get("slide_number")
                if slide_num:
                    if slide_num not in visuals_by_slide:
                        visuals_by_slide[slide_num] = []
                    visuals_by_slide[slide_num].append(visual_data)

            chunks = []
            for slide_number, slide in enumerate(prs.slides, start=1):
                # Extract text from this slide's shapes
                slide_text_parts = [
                    shape.text
                    for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text
                ]
                slide_text = " ".join(slide_text_parts).strip()

                if slide_text:
                    # Create a single chunk per slide with the slide text
                    text_chunk = {
                        "content": slide_text,
                        "tokens": self.num_tokens_from_string(slide_text),
                        "location": {"slide_number": slide_number},
                        "canSplit": True,
                    }
                    chunks.append(text_chunk)

                # Process any visuals that belong to this slide (O(1) lookup)
                slide_visuals = visuals_by_slide.get(slide_number, [])
                for visual_data in slide_visuals:
                    # Check if visual has been processed (has transcription)
                    if visual_data.get("transcription"):
                        # Create visual chunk using the format function
                        visual_chunk = format_visual_chunk_data(visual_data, self.num_tokens_from_string)
                        chunks.append(visual_chunk)

            return chunks

    ### Visual Data Extraction ###
    def preprocess_pptx_visuals(self, file_content):
        """
        Extract visual content and create a visual map for LLM processing.
        Modifies the PowerPoint file to include visual markers that MarkItDown will preserve.
        """

        # First pass: Extract visual data with images
        visual_map = {}
        shape_to_marker = {}  # Track which shapes get which markers
        
        with io.BytesIO(file_content) as f:
            presentation_with_images = Presentation(f)

        # Process each slide to extract visuals and track shape relationships
        for slide_idx, slide in enumerate(presentation_with_images.slides, start=1):
            slide_visuals = self.extract_slide_visuals(slide, slide_idx)
            visual_map.update(slide_visuals)
            
            # Track which shapes correspond to which markers
            for marker, visual_data in slide_visuals.items():
                # We need to identify shapes - use a combination of slide + shape index
                shape_id = f"slide_{slide_idx}_shape_{visual_data.get('shape_index', 0)}"
                shape_to_marker[shape_id] = marker

        # Second pass: Create clean presentation with markers instead of images
        with io.BytesIO(file_content) as f:
            clean_presentation = Presentation(f)
            
        # Inject markers into clean presentation
        self._inject_markers_into_presentation(clean_presentation, shape_to_marker)

        # Convert modified presentation to bytes
        modified_content = self._presentation_to_bytes(clean_presentation)
        return modified_content, visual_map

    def _inject_markers_into_presentation(self, presentation, shape_to_marker):
        """
        Replace visual shapes with text boxes containing markers.
        """
        for slide_idx, slide in enumerate(presentation.slides, start=1):
            # Build list of shapes to process (avoid modifying collection while iterating)
            shapes_to_process = []
            
            for shape_index, shape in enumerate(slide.shapes):
                is_visual, visual_type = self.is_extractable_visual(shape)
                
                if is_visual:
                    # Generate the same shape_id used during extraction
                    shape_id = f"slide_{slide_idx}_shape_{shape_index}"
                    marker = shape_to_marker.get(shape_id)
                    
                    if marker:
                        shapes_to_process.append({
                            'shape': shape,
                            'marker': marker,
                            'left': shape.left,
                            'top': shape.top,
                            'width': shape.width,
                            'height': shape.height
                        })
            
            # Process shapes (remove and add markers)
            for shape_info in shapes_to_process:
                # Remove the original visual shape
                self._remove_shape_from_slide(slide, shape_info['shape'])
                
                # Add marker text box in its place
                self._add_marker_text_box(slide, {
                    'marker': shape_info['marker'],
                    'left': shape_info['left'],
                    'top': shape_info['top'],
                    'width': shape_info['width'],
                    'height': shape_info['height']
                })

    def _remove_shape_from_slide(self, slide, shape):
        """
        Remove a shape from a slide.
        """
        try:
            # Get the slide's shape collection
            shapes = slide.shapes
            
            # Find the shape in the collection and remove it
            # We need to work with the underlying XML elements
            shape_element = shape._element
            parent = shape_element.getparent()
            
            if parent is not None:
                parent.remove(shape_element)
                print(f"Successfully removed shape: {getattr(shape, 'name', 'unnamed')}")
            else:
                print(f"Warning: Could not find parent for shape: {getattr(shape, 'name', 'unnamed')}")
                
        except Exception as e:
            print(f"Error removing shape: {e}")
            # Try alternative removal method
            try:
                # Alternative approach: remove from spTree directly
                slide._slide.spTree.remove(shape._element)
                print(f"Successfully removed shape using alternative method")
            except Exception as e2:
                print(f"Alternative removal also failed: {e2}")

    def _add_marker_text_box(self, slide, marker_info):
        """
        Add a text box with the visual marker to the slide.
        """
        try:
            # Add a text box at the same position as the original shape
            text_box = slide.shapes.add_textbox(
                marker_info['left'],
                marker_info['top'], 
                marker_info['width'],
                marker_info['height']
            )
            
            # Set the text to the marker
            text_frame = text_box.text_frame
            text_frame.text = marker_info['marker']
            
            # Make the text visible but not too prominent
            paragraph = text_frame.paragraphs[0]
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = RGBColor(128, 128, 128)  # Gray color
            
        except Exception as e:
            print(f"Warning: Could not add marker text box: {e}")

    def _presentation_to_bytes(self, presentation):
        """
        Convert a presentation object back to bytes.
        """
        output_buffer = io.BytesIO()
        presentation.save(output_buffer)
        return output_buffer.getvalue()

    def is_extractable_visual(self, shape):
        """Determine if a shape contains extractable visual content"""
        # ALWAYS extract: Images (photos, screenshots, etc.)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:  # 13
            return True, VisualType.IMAGE.value

        # EXTRACT: Complex diagrams that markitdown might miss
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:  # 6
            if self.has_significant_content(shape):
                return True, VisualType.DIAGRAM.value

        # SMART EXTRACT: FREEFORM shapes (custom drawn) - only if they're likely visual content
        if shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:  # 5
            if self.is_meaningful_freeform(shape):
                return True, VisualType.FREEFORM.value

        # SKIP: AUTO_SHAPE (text boxes, basic shapes) - MarkItDown handles these well
        # SKIP charts - markitdown handles these perfectly
        # SKIP: Text boxes, connectors, small decorative elements
        return False, None

    def has_significant_content(self, group_shape):
        """Check if grouped shape contains meaningful visual content"""
        child_count = len(list(group_shape.shapes))

        # Groups with multiple elements are likely meaningful diagrams
        if child_count >= 3:
            return True

        # Check if any child has substantial size
        for child in group_shape.shapes:
            if self.is_significant_size(child):
                return True

        return False

    def is_significant_size(self, shape):
        """Determine if shape is large enough to contain meaningful content"""
        # Convert EMU (English Metric Units) to pixels approximately
        width_px = shape.width / 9525  # Rough conversion
        height_px = shape.height / 9525

        # Minimum size thresholds
        MIN_WIDTH = 100  # pixels
        MIN_HEIGHT = 50  # pixels
        MIN_AREA = 8000  # square pixels

        return (
            width_px >= MIN_WIDTH
            and height_px >= MIN_HEIGHT
            and width_px * height_px >= MIN_AREA
        )

    def is_meaningful_freeform(self, shape):
        """Determine if a freeform shape contains meaningful visual content vs just text"""

        # Must be significant size to be meaningful
        if not self.is_significant_size(shape):
            return False

        # Get text content if any
        text_content = ""
        try:
            if hasattr(shape, "text") and shape.text:
                text_content = shape.text.strip()
        except:
            pass

        # If it has no text, it's likely a pure visual element (arrow, custom shape, etc.)
        if not text_content:
            return True

        # If it has text, check if it's likely just a fancy text container
        # vs. a meaningful diagram with text labels

        # Short text in large shape = likely meaningful visual with labels
        if len(text_content) < 50 and self.is_very_large(shape):
            return True

        # Very long text = probably just a fancy text box
        if len(text_content) > 200:
            return False

        # Medium text in medium shape = extract it (could be callout, annotation, etc.)
        return True

    def is_very_large(self, shape):
        """Check if shape is very large (likely a diagram vs text container)"""
        width_px = shape.width / 9525
        height_px = shape.height / 9525
        area = width_px * height_px

        # Very large shapes are likely diagrams, not just text
        return area > 50000  # Much larger than our normal threshold

    def extract_slide_visuals(self, slide, slide_number):
        """Extract all visual content from a single slide"""

        visuals = {}
        # Only count the visual types we actually extract
        type_counters = {
            VisualType.IMAGE.value: 0,
            VisualType.DIAGRAM.value: 0,
            VisualType.FREEFORM.value: 0,
        }

        for shape_index, shape in enumerate(slide.shapes):
            is_visual, visual_type = self.is_extractable_visual(shape)

            if is_visual:
                type_counters[visual_type] += 1
                marker = f"<Visual#{slide_number}_{visual_type}_{type_counters[visual_type]}>"

                visual_data = self.extract_visual_data(shape, visual_type, slide_number)

                # Only add to map if extraction was successful
                if visual_data is not None:
                    # Add shape index for tracking
                    visual_data['shape_index'] = shape_index
                    visuals[marker] = visual_data
                else:
                    print(f"Skipping visual {marker} due to extraction failure")

        return visuals

    def extract_visual_data(self, shape, visual_type, slide_number):
        """Extract the actual visual data from a shape"""
        print(f"Extracting visual data for {visual_type} on slide {slide_number}")

        if visual_type == VisualType.IMAGE.value:
            return self.extract_image_data(shape, slide_number)
        elif visual_type in [VisualType.DIAGRAM.value, VisualType.FREEFORM.value]:
            return self.extract_shape_data(shape, visual_type, slide_number)

    def extract_image_data(self, shape, slide_number):
        """Extract image data and metadata"""

        # Extract the image using the correct python-pptx approach
        image_bytes = None
        original_format = None

        try:
            slide_part = shape.part
            rId = shape._element.blip_rId

            # Use rels to get the relationship, then target_part
            if hasattr(slide_part, "rels") and rId in slide_part.rels:
                rel = slide_part.rels[rId]
                image_part = rel.target_part
                image_bytes = image_part.blob
                original_format = image_part.content_type
            else:
                print(f"Could not find image relationship {rId} on slide {slide_number}" )
                return None

        except Exception as e:
            print(f"Image extraction failed on slide {slide_number}: {e}")
            return None

        # Convert format if necessary
        final_image_bytes, final_format = ensure_supported_format(
            image_bytes, original_format
        )

        # Generate hash for deduplication
        content_hash = hash_visual_data(final_image_bytes)

        # Extract additional metadata
        metadata = self.extract_image_metadata(shape)

        return self._create_visual_data_structure(
            VisualType.IMAGE.value,
            final_format,
            final_image_bytes,
            content_hash,
            {"slide_number": slide_number},
            metadata,
            original_format,
        )

    def extract_image_metadata(self, shape):
        """Extract metadata from image shape"""
        metadata = {}

        try:
            # Alt text (description)
            if hasattr(shape, "_element") and hasattr(shape._element, "get"):
                # Try to get alt text from various PowerPoint properties
                alt_text = ""

                # Method 1: Try descr attribute (common alt text location)
                try:
                    if hasattr(shape._element, "nvPicPr"):
                        cNvPr = shape._element.nvPicPr.cNvPr
                        if cNvPr.get("descr"):
                            alt_text = cNvPr.get("descr")
                except:
                    pass

                # Method 2: Try name attribute as fallback
                if not alt_text:
                    try:
                        if hasattr(shape._element, "nvPicPr"):
                            cNvPr = shape._element.nvPicPr.cNvPr
                            if cNvPr.get("name"):
                                alt_text = cNvPr.get("name")
                    except:
                        pass

                metadata["alt_text"] = alt_text

            # Shape name
            metadata["shape_name"] = getattr(shape, "name", "")

            # Title (from text content if available)
            metadata["title"] = ""
            if hasattr(shape, "text_frame") and shape.text_frame:
                metadata["title"] = shape.text_frame.text.strip()
            elif hasattr(shape, "text") and shape.text:
                metadata["title"] = shape.text.strip()

            # Hyperlink
            metadata["hyperlink"] = ""
            try:
                if (
                    hasattr(shape, "click_action")
                    and hasattr(shape.click_action, "hyperlink")
                    and shape.click_action.hyperlink
                    and hasattr(shape.click_action.hyperlink, "address")
                ):
                    metadata["hyperlink"] = shape.click_action.hyperlink.address or ""
            except Exception:
                # Group shapes and some other shape types don't support click actions
                pass

        except Exception as e:
            print(f"Error extracting image metadata: {e}")

        return metadata

    def render_shape_as_image(self, shape, label="Visual"):
        """
        Render a shape as PNG image for Lambda environment.
        Creates a simple placeholder image with shape metadata.
        """
        try:
            # Get shape dimensions in pixels (approximate)
            width_px = max(200, min(800, int(shape.width / 9525)))
            height_px = max(150, min(600, int(shape.height / 9525)))

            # Create a new image with white background
            img = Image.new("RGB", (width_px, height_px), color="white")
            draw = ImageDraw.Draw(img)

            # Try to use a default font
            try:
                font = ImageFont.load_default()
            except:
                font = None

            # Draw a border
            draw.rectangle(
                [(0, 0), (width_px - 1, height_px - 1)], outline="black", width=2
            )

            # Add label text
            if font:
                # Calculate text position (centered)
                bbox = draw.textbbox((0, 0), label, font=font)
                text_width = bbox[2] - bbox[0]
                text_height = bbox[3] - bbox[1]
                x = (width_px - text_width) // 2
                y = (height_px - text_height) // 2
                draw.text((x, y), label, fill="black", font=font)

            # Add shape type info
            shape_info = f"Type: {shape.shape_type}"
            if font and height_px > 100:
                bbox = draw.textbbox((0, 0), shape_info, font=font)
                text_width = bbox[2] - bbox[0]
                x = (width_px - text_width) // 2
                y = y + 30
                draw.text((x, y), shape_info, fill="gray", font=font)

            # Save to bytes
            img_bytes = io.BytesIO()
            img.save(img_bytes, format="PNG")
            return img_bytes.getvalue()

        except Exception as e:
            print(f"Error rendering shape as image: {e}")
            # Create minimal fallback image
            img = Image.new("RGB", (200, 150), color="lightgray")
            draw = ImageDraw.Draw(img)
            draw.rectangle([(0, 0), (199, 149)], outline="black", width=1)
            draw.text((10, 70), "Visual Content", fill="black")

            img_bytes = io.BytesIO()
            img.save(img_bytes, format="PNG")
            return img_bytes.getvalue()

    def extract_shape_data(self, shape, visual_type, slide_number):
        """Extract shape/diagram data"""

        # For complex shapes, we'll render them as images
        shape_image = self.render_shape_as_image(
            shape, f"{visual_type.title()} Diagram"
        )

        # Generate hash for deduplication
        content_hash = hash_visual_data(shape_image)

        # Try to get any text content from the shape
        shape_text = ""
        if hasattr(shape, "text") and shape.text:
            shape_text = shape.text.strip()

        # Extract shape metadata (alt text, hyperlinks, etc.)
        metadata = self.extract_shape_metadata(shape)

        return self._create_visual_data_structure(
            visual_type,
            PNG,
            shape_image,
            content_hash,
            shape_text,
            {"slide_number": slide_number},
            metadata,
            "",
        )

    def extract_shape_metadata(self, shape):
        """Extract metadata from any shape (charts, shapes, etc.)"""
        metadata = {}

        try:
            # Alt text (description) - works for all shape types
            alt_text = ""
            try:
                if hasattr(shape, "_element"):
                    # For charts and other shapes, try different XML paths
                    element = shape._element

                    # Try chart-specific path
                    if hasattr(element, "nvGraphicFramePr"):
                        cNvPr = element.nvGraphicFramePr.cNvPr
                        if cNvPr.get("descr"):
                            alt_text = cNvPr.get("descr")

                    # Try general shape path
                    elif hasattr(element, "nvSpPr"):
                        cNvPr = element.nvSpPr.cNvPr
                        if cNvPr.get("descr"):
                            alt_text = cNvPr.get("descr")

                    # Fallback: try name attribute
                    if not alt_text and hasattr(element, "nvSpPr"):
                        cNvPr = element.nvSpPr.cNvPr
                        if cNvPr.get("name"):
                            alt_text = cNvPr.get("name")
            except:
                pass

            metadata["alt_text"] = alt_text

            # Shape name
            metadata["shape_name"] = getattr(shape, "name", "")

            # Title (from text content if available)
            metadata["title"] = ""
            if hasattr(shape, "text_frame") and shape.text_frame:
                metadata["title"] = shape.text_frame.text.strip()
            elif hasattr(shape, "text") and shape.text:
                metadata["title"] = shape.text.strip()

            # Hyperlink
            metadata["hyperlink"] = ""
            try:
                if (
                    hasattr(shape, "click_action")
                    and hasattr(shape.click_action, "hyperlink")
                    and shape.click_action.hyperlink
                    and hasattr(shape.click_action.hyperlink, "address")
                ):
                    metadata["hyperlink"] = shape.click_action.hyperlink.address or ""
            except Exception:
                # Group shapes and some other shape types don't support click actions
                pass

        except Exception as e:
            print(f"Error extracting shape metadata: {e}")

        return metadata

    def _create_visual_data_structure(
        self, type, format, data, hash, location, metadata, original_format
    ):
        return {
            "type": type,
            "format": format,
            "data": data,
            "hash": hash,
            "location": location,
            "alt_text": metadata.get("alt_text", ""),
            "title": metadata.get("title", ""),
            "hyperlink": metadata.get("hyperlink", ""),
            "shape_name": metadata.get("shape_name", ""),
            "original_format": original_format,
        }
